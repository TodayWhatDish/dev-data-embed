"""오늘뭐먹냥 포트폴리오 발표자료(pptx) 생성 스크립트.

python -m docs.deck.build_pptx  (또는) python docs/deck/build_pptx.py
docs/deck/PT.md 에 정리된 주제를 바탕으로, 저장소 코드/DB/pytest 실행 결과에서
직접 확인한 수치만 근거로 사용한다 (2026-09-04 기준).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE

# ── 팔레트 (dataviz 스킬 검증 팔레트: 카테고리 순서 고정) ─────────────────
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SECOND = RGBColor(0x52, 0x51, 0x4E)
BG = RGBColor(0xFC, 0xFC, 0xFB)
ORANGE = RGBColor(0xEB, 0x68, 0x34)   # 브랜드 강조색 (slot 2)
BLUE = RGBColor(0x2A, 0x78, 0xD6)     # 구조/기술 강조색 (slot 1)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)     # slot 3
CARD_BG = RGBColor(0xF3, 0xF1, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # 뒤로 보내기(가장 아래로)
    bg._element.getparent().remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def _set_font(run, size, color=INK, bold=False, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        _set_font(r, size, color, bold)
    return box


def add_bullets(slide, x, y, w, h, items, size=16, color=INK, gap=6, bold_lead=False):
    """items: list[str] 또는 (lead, rest) 튜플. '•' 를 직접 붙인다."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.2
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = f"•  {lead}"
            _set_font(r1, size, color, True)
            if rest:
                r2 = p.add_run()
                r2.text = "  " + rest
                _set_font(r2, size - 1, INK_SECOND, False)
        else:
            r = p.add_run()
            r.text = f"•  {item}"
            _set_font(r, size, color, bold_lead)
    return box


def add_kicker_title(slide, kicker, title, num):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(6), Inches(0.4),
              kicker, size=14, color=ORANGE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.68), Inches(11.5), Inches(0.9),
              title, size=26, color=INK, bold=True)
    add_text(slide, Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.35),
              f"{num:02d}", size=12, color=INK_SECOND, align=PP_ALIGN.RIGHT)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.18), Inches(1.0), Pt(3))
    rule.fill.solid(); rule.fill.fore_color.rgb = ORANGE; rule.line.fill.background()
    rule.shadow.inherit = False


def add_card(slide, x, y, w, h, title, body, accent=ORANGE, title_size=15, body_size=13):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.06
    box.fill.solid(); box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = accent; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.12)
    p0 = tf.paragraphs[0]
    r0 = p0.add_run(); r0.text = title
    _set_font(r0, title_size, accent, True)
    p1 = tf.add_paragraph()
    p1.space_before = Pt(4)
    p1.line_spacing = 1.15
    r1 = p1.add_run(); r1.text = body
    _set_font(r1, body_size, INK)
    return box


def add_flow(slide, x, y, w, h, steps, accent=BLUE, size=13):
    """steps: list[str]. 박스 n개를 화살표로 잇는다."""
    n = len(steps)
    gap = Inches(0.35)
    box_w = Emu(int((w - gap * (n - 1)) / n))
    for i, step in enumerate(steps):
        bx = x + i * (box_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, box_w, h)
        box.adjustments[0] = 0.12
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = accent; box.line.width = Pt(1.25)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = step
        _set_font(r, size, INK, True)
        if i < n - 1:
            ax = bx + box_w
            add_text(slide, ax, y, gap, h, "→", size=18, color=accent,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_table(slide, x, y, w, h, header, rows, col_widths=None, size=12):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(header), x, y, w, h)
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htext in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = ORANGE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = htext
        _set_font(r, size, WHITE, True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else CARD_BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            _set_font(r, size, INK, j == 0)
    return table


def add_bar_chart(slide, x, y, w, h, categories, series, colors, title=None):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series.items():
        data.add_series(name, values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data)
    chart = gframe.chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = FONT
    for i, plot in enumerate(chart.plots):
        for s, color in zip(plot.series, colors):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = color
    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(11); cat_ax.tick_labels.font.name = FONT
    val_ax = chart.value_axis
    val_ax.tick_labels.font.size = Pt(10); val_ax.tick_labels.font.name = FONT
    val_ax.has_major_gridlines = False
    if title:
        add_text(slide, x, y - Inches(0.35), w, Inches(0.3), title, size=13,
                  color=INK_SECOND, bold=True)
    return chart


def add_placeholder(slide, x, y, w, h, label):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.04
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xE6, 0xE2)
    box.line.color.rgb = INK_SECOND; box.line.width = Pt(1)
    box.line.dash_style = 3  # dash
    box.shadow.inherit = False
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    _set_font(r, 13, INK_SECOND, False, True)


# ══════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    page = [0]

    def num():
        page[0] += 1
        return page[0]

    # ── 1. 표지 ──────────────────────────────────────────────────────
    s = new_slide(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.35), SLIDE_W, Inches(2.5))
    band.fill.solid(); band.fill.fore_color.rgb = ORANGE
    band.line.fill.background(); band.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.2),
              "오늘뭐먹냥", size=54, color=WHITE, bold=True)
    add_text(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.7),
              "우리 아이의 소중한 한입", size=22, color=WHITE)
    add_text(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.5),
              "근거 있는 리뷰 기반 반려동물 사료·간식 추천 서비스 (RAG)", size=16, color=INK_SECOND)
    add_text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
              "허재성 · 박승호 · 최우진 · 한인혁 · 최종명", size=14, color=INK)

    # ── 2. 목차 ──────────────────────────────────────────────────────
    s = new_slide(prs)
    add_kicker_title(s, "AGENDA", "목차", num())
    agenda = [
        "서비스 선정 이유와 비즈니스 모델",
        "데이터 설계 — 스키마와 계층 구조",
        "개인정보 마스킹과 품질 검증",
        "관리자 대시보드와 LLM 반증(Fact-check)",
        "리뷰 데이터 변동과 임베딩 자동 갱신",
        "FastAPI와 계층형 아키텍처",
        "LangChain과 상용 LLM·임베딩 모델 선정",
        "서비스 화면",
    ]
    left = agenda[:4]; right = agenda[4:]
    for col, items in [(0, left), (1, right)]:
        x = Inches(0.7 + col * 6.1)
        tf_items = [(f"{i+1+col*4:02d}", t) for i, t in enumerate(items)]
        add_bullets(s, x, Inches(1.7), Inches(5.6), Inches(4.8), tf_items, size=17, gap=18)

    # ── 3. 서비스 선정 이유 & BM ─────────────────────────────────────
    s = new_slide(prs)
    add_kicker_title(s, "01 · 서비스 선정 이유", "돈이 되는 도메인인가 — 왜 반려동물인가", num())
    add_card(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(2.6),
              "도메인이 깊다",
              "개·고양이로 한정해도 품종·체중·체형(BCS 5단계)·중성화 여부·알러지·급여목적까지 "
              "정형 속성만 8종 이상. 축종을 하나 늘리는 것도 품종표·체형척도·급여기준이 통째로 "
              "따라와야 해서 '행 하나 추가'가 안 되는 도메인.", accent=BLUE)
    add_card(s, Inches(0.6), Inches(4.3), Inches(6.0), Inches(2.6),
              "경쟁 구도 대비 차별점",
              "핏펫(AI 진단+자사몰 커머스+구독), 펫스튜어드(무료 AI 진단+네이버 스토어 수수료) 모두 "
              "'얼마나 많은 리뷰를 분석했나'로 경쟁. 오늘뭐먹냥은 '이 리뷰가 정말 우리 아이 이야기인가' — "
              "근거(리뷰 원문)를 남겨 추천 품질을 사후에 검증할 수 있게 함.", accent=ORANGE)
    add_table(s, Inches(6.9), Inches(1.5), Inches(5.85), Inches(2.6),
               ["항목", "내용"],
               [["사업 방식", "B2C"],
                ["타깃", "반려동물 보호자 / 용품 판매업체"],
                ["수익 모델", "프리미엄 AI 기능 구독"],
                ["과금 구조", "7일 무료 체험 → 월정액"]],
               col_widths=[Inches(1.9), Inches(3.95)])
    add_card(s, Inches(6.9), Inches(4.3), Inches(5.85), Inches(2.6),
              "추천의 첫 관문 = 정형 필터",
              "축종이 뒤섞이면 그 뒤의 어떤 랭킹도 무의미하므로, SQL이 축종·알러지를 먼저 거르고 "
              "LLM은 이미 걸러진 후보 위에서만 판단한다 (오탐 시 사고 비용이 큰 도메인이기 때문).",
              accent=AQUA)

    # ── 4. 데이터 설계 ────────────────────────────────────────────────
    s = new_slide(prs)
    add_kicker_title(s, "02 · 데이터 설계", "스키마 설계와 계층 구조", num())
    add_text(s, Inches(0.6), Inches(1.55), Inches(11.8), Inches(0.35),
              "의존 방향은 위 → 아래 한 방향. repositories 는 domain 을 모르고, domain 은 SQL 을 모른다.",
              size=14, color=INK_SECOND)
    add_flow(s, Inches(0.6), Inches(2.1), Inches(11.8), Inches(0.75),
              ["app/api  (HTTP만)", "app/features  (유스케이스 조립)",
               "domain / repositories  (계산 · SQL)", "core/db.py  (SQLite 유일 진입점)"],
              accent=BLUE, size=12)
    add_card(s, Inches(0.6), Inches(3.25), Inches(5.85), Inches(1.9),
              "다중값 속성은 별도 N:M 테이블로",
              "product_ingredient / product_animal_category / product_feeding_purpose / "
              "ingredient_allergen / pet_allergy — 콤마 문자열이나 JSON 한 컬럼으로 접지 않는다. "
              "마스터 데이터는 기동 시 한 번 읽어 싱글턴 캐시에 올려, 조인 이름 규칙이 SQL에 흩어지지 않게 한다.",
              accent=ORANGE, body_size=12.5)
    add_card(s, Inches(6.65), Inches(3.25), Inches(6.1), Inches(1.9),
              "계층별 책임 분리",
              "repositories = SELECT 전담(스키마 변경 대응 지점 하나로 좁힘) · domain = 마스터 캐시·판정 로직"
              "(DB 없이 자체검증 가능) · features = 유스케이스 순서(선필터→랭킹→조립).",
              accent=BLUE, body_size=12.5)
    add_table(s, Inches(0.6), Inches(5.35), Inches(12.15), Inches(1.55),
               ["실측 규모 (data/pet_reco.db, 2026-09-04)", "테이블", "뷰", "상품", "리뷰", "반려동물"],
               [["운영 DB", "22개", "2개", "200개", "2,001건", "369마리"]],
               col_widths=[Inches(4.15)] + [Inches(1.6)] * 5)

    # ── 5. 마스킹 & 품질 검증 ────────────────────────────────────────
    s = new_slide(prs)
    add_kicker_title(s, "03 · 데이터 마스킹과 평가", "완벽 차단이 아니라 확실한 것만 — 마스킹 정책과 검증", num())
    add_bullets(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(2.5), [
        ("원본은 보존", "DB 저장값은 안 건드림 — 마스킹은 되돌릴 수 없는 파괴적 변환이라 밖(LLM)으로 나가는 글에서만 적용"),
        ("완벽 차단 대신 확실한 것만", "우회 수법을 다 쫓으면 규칙이 조여져 평범한 후기가 먼저 지워짐 — '새는 쪽'을 택함"),
        ("문장째가 아니라 그 부분만 치환", "마침표 없는 글에서 문장 단위로 지우면 상품 설명까지 통째로 날아가는 오탐 방지"),
    ], size=13.5, gap=10)
    add_table(s, Inches(6.9), Inches(1.55), Inches(5.85), Inches(2.5),
               ["입력 → 출력 (실제 assert)", ""],
               [["'구매는 010-1234-5678 로'", "→ [전화번호]"],
                ["'O1O-1234-5678 이거로' (우회)", "→ [전화번호] 탐지"],
                ["'배송 문의 연락했더니...'", "그대로 보존 (오탐 방지)"],
                ["'...남기더라고요' (후기 106/2000건)", "그대로 보존 (규칙 제외)"]],
               col_widths=[Inches(4.35), Inches(1.5)], size=11)
    add_bar_chart(s, Inches(0.6), Inches(4.7), Inches(5.7), Inches(2.15),
                   ["규칙 카테고리 수"], {"개수": [8]}, [ORANGE],
                   title="마스킹 정규식 규칙 (app/domain/masking.py: 전화/한글전화/이메일/카톡/기타채널/SNS/주소/이름)")
    add_card(s, Inches(6.55), Inches(4.7), Inches(6.2), Inches(2.15),
              f"pytest 자체검증 — 2026-09-04 기준 실행 결과",
              "SELFCHECKS 12개 모듈(마스킹·알러지·품종·청킹 연동 등) 파라미터라이즈 → 28 passed, 1 deselected(느린 임베딩 로드 테스트). "
              "마스킹 모듈만도 40개 assert(반복 케이스 포함 70+ 건) 전원 통과.", accent=AQUA, body_size=12.5)

    # ── 6. 관리자 대시보드 & 반증 ─────────────────────────────────────
    s = new_slide(prs)
    add_kicker_title(s, "04 · 관리자 대시보드", "LLM이 LLM을 반증(Fact-check)한다", num())
    add_flow(s, Inches(0.6), Inches(1.6), Inches(11.8), Inches(0.8),
              ["관리자 질문", "답변 생성 LLM\n(claude-sonnet-5)", "반증 LLM\n(claude-haiku-4-5)",
               "NDJSON 스트리밍\n(관리자 화면)"], accent=ORANGE, size=12)
    add_card(s, Inches(0.6), Inches(2.75), Inches(6.0), Inches(2.2),
              "왜 같은 모델이 스스로 채점하지 않나",
              "같은 호출·같은 모델이 스스로 채점하면 자기가 만든 답을 정당화하는 편향이 생긴다. "
              "answering.verify() 는 별도의 chat.with_structured_output(FactCheck) 호출로 "
              "[고객 정보]라는 사실 자료와만 대조한다. (app/features/answering.py:22-26)",
              accent=BLUE, body_size=12.5)
    add_card(s, Inches(6.9), Inches(2.75), Inches(5.85), Inches(2.2),
              "왜 타사가 아니라 Anthropic 계열끼리인가",
              "편향 제거에 필요한 조건은 '다른 벤더'가 아니라 '다른 모델·다른 독립 호출'. "
              "생성은 Sonnet, 반증은 저비용·저지연 Haiku로 나눠, 벤더를 늘리지 않고도 같은 요구를 "
              "더 싼 비용으로 만족시킨다.", accent=AQUA, body_size=12.5)
    add_card(s, Inches(0.6), Inches(5.15), Inches(12.15), Inches(1.55),
              "생성과 반증은 독립적으로 실패한다",
              "app/api/routes/ask.py:41-63 의 generate() 가 답변 생성과 반증을 각각 try/except로 감싼다 — "
              "반증이 실패해도 이미 스트리밍된 답변은 사라지지 않는다.", accent=ORANGE, body_size=13)

    # ── 7. 임베딩 자동 갱신 ──────────────────────────────────────────
    s = new_slide(prs)
    add_kicker_title(s, "05 · 데이터 변동 대응", "새 리뷰가 쌓일 때, 전체가 아니라 변경분만 재색인", num())
    add_flow(s, Inches(0.6), Inches(1.6), Inches(11.8), Inches(0.8),
              ["새 리뷰 INSERT", "pipeline.chunk\n(조각 +1)", "pipeline.embed\n(새로 1 / 그대로 N / 지움 0)",
               "chunk_vectors\n+ embedding_meta"], accent=BLUE, size=12)
    add_card(s, Inches(0.6), Inches(2.75), Inches(6.0), Inches(2.5),
              "재색인 감지 — check_freshness",
              "embedding_meta 테이블에 벡터가 '언제·무엇(모델·버전)으로' 만들어졌는지 기록해 두고, "
              "이 기록과 현재 청크를 비교해 바뀐 것만 다시 임베딩한다. 가장 흔한 사고는 데이터를 바꾸고 "
              "재색인을 잊는 것 — 이 체크가 그 사고를 잡는다.", accent=ORANGE, body_size=12.5)
    add_card(s, Inches(6.9), Inches(2.75), Inches(5.85), Inches(2.5),
              "증분 검증 — tests/incremental_embed",
              "CSV 전체 재적재로는 '딱 한 건 늘었다'를 볼 수 없어, 회원가입과 동일한 INSERT 경로로 "
              "리뷰 1건만 직접 넣고 chunk/embed를 돌려 '새로 1 / 그대로 N / 지움 0'을 그대로 확인한다. "
              "정리(cleanup) 후에는 반대로 '지움 1'까지 검증.", accent=AQUA, body_size=12.5)
    add_text(s, Inches(0.6), Inches(5.5), Inches(12.15), Inches(1.2),
              "→ 서비스 운영 중 새 구매·리뷰가 쌓여도 전체 파이프라인을 다시 돌릴 필요 없이, "
              "chunk → embed 두 명령만으로 벡터 인덱스가 최신 상태를 유지한다.",
              size=14, color=INK_SECOND)

    # ── 8. FastAPI 선정 이유 ────────────────────────────────────────
    s = new_slide(prs)
    add_kicker_title(s, "06 · 서버 아키텍처", "Flask·Django 대신 FastAPI를 고른 이유", num())
    add_table(s, Inches(0.6), Inches(1.55), Inches(12.15), Inches(2.15),
               ["근거", "코드 위치"],
               [["Pydantic 스키마로 요청/응답을 자동 검증 — 라우터는 HTTP↔값 변환만 하고 SQL·비즈니스 로직은 모른다",
                 "app/api/schemas.py"],
                ["lifespan 하나로 서버 기동 시 벡터 DB 커넥션을 열고, 종료 시 닫는다 (배치·서버가 같은 검색 함수 공유)",
                 "app/api/lifespan.py"],
                ["동기 엔드포인트도 스레드풀에서 돌아, SQLite 스레드로컬 커넥션(WAL) 모델과 자연스럽게 맞물린다",
                 "app/core/db.py"],
                ["자동 생성되는 OpenAPI 스펙이 dev-web(별도 레포) 프런트와의 계약 문서 역할을 대신한다",
                 "/docs"]],
               col_widths=[Inches(9.6), Inches(2.55)], size=12)
    add_flow(s, Inches(0.6), Inches(4.1), Inches(12.15), Inches(0.75),
              ["api  (HTTP만)", "features  (유스케이스)", "domain / repositories  (계산 · SQL)",
               "core/db.py"], accent=ORANGE, size=12)
    add_card(s, Inches(0.6), Inches(5.2), Inches(12.15), Inches(1.5),
              "요청 한 건의 여정",
              "/ask 라우트는 답변 생성과 반증 두 번의 LLM 호출을 NDJSON으로 순서대로 흘려보낸다. "
              "main.py는 라우터 등록과 정적 파일(/static/admin) 마운트만 하고, DB 연결·예외 변환은 "
              "각 계층에 위임한다 — main.py는 '조립'만 한다.", accent=BLUE, body_size=13)

    # ── 9. LangChain & 모델 선정 ─────────────────────────────────────
    s = new_slide(prs)
    add_kicker_title(s, "07 · LangChain과 모델 선정", "프로바이더 스위치 하나로 로컬↔상용을 오간다", num())
    add_card(s, Inches(0.6), Inches(1.55), Inches(12.15), Inches(1.55),
              "LangChain의 역할 — init_chat_model()",
              "LLM_PROVIDER 값만 보고 ChatAnthropic/ChatOpenAI 등 알맞은 클라이언트를 골라준다. "
              "분기는 core/config.py 한 곳에만 있고 features/*는 프로바이더가 뭐든 코드가 안 바뀐다 — "
              "개발 중엔 provider=openai + base_url만 로컬 Ollama로 돌려 비용 없이 검증, 운영은 anthropic. "
              "(app/adapters/stores/llm.py)", accent=BLUE, body_size=12.5)
    add_bar_chart(s, Inches(0.6), Inches(3.7), Inches(6.6), Inches(3.1),
                   ["recall@1", "recall@3", "recall@10"],
                   {"e5-small (384d, 14ms)": [8.7, 18.9, 34.9],
                    "e5-base (768d, 41ms)": [10.2, 17.8, 30.5],
                    "bge-m3 (1024d, 132ms)": [6.9, 18.9, 36.4]},
                   [BLUE, ORANGE, AQUA],
                   title="임베딩 모델 비교 (data/eval, 홀드아웃 275건, 단위 %)")
    add_bullets(s, Inches(7.5), Inches(3.75), Inches(5.25), Inches(3.0), [
        ("표본 66건 기준 잡음 폭 ±1.5%p", "noise_band()가 2,000회 복원추출로 이 폭을 먼저 잰다 — 이 안의 차이는 모델 우열이 아니라 표본 운"),
        ("속도·비용 축에서 결정", "recall이 잡음 폭 안이라면 쿼리 지연(14ms~132ms)과 차원(비용)으로 고른다"),
        ("Anthropic 선택 이유", "구조화 출력(with_structured_output) 안정성 + 반증용 저비용 Haiku까지 한 벤더 안에서 확보"),
    ], size=12.5, gap=10)

    # ── 10. 서비스 화면 ──────────────────────────────────────────────
    s = new_slide(prs)
    add_kicker_title(s, "08 · 서비스 화면", "데모", num())
    add_placeholder(s, Inches(0.6), Inches(1.6), Inches(5.85), Inches(5.1),
                      "[스크린샷: 관리자 대시보드\n(질문 스트리밍 · 반증 · 마스킹)]")
    add_placeholder(s, Inches(6.65), Inches(1.6), Inches(6.1), Inches(5.1),
                      "[스크린샷: 고객 페이지\n추천 결과 화면]")

    # ── 11. 마무리 ───────────────────────────────────────────────────
    s = new_slide(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    band.fill.solid(); band.fill.fore_color.rgb = ORANGE
    band.line.fill.background(); band.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.0),
              '"우리 아이를 아는 AI"', size=36, color=WHITE, bold=True)
    add_bullets(s, Inches(0.9), Inches(3.3), Inches(11.0), Inches(2.5), [
        "프로필 개인화 — 축종·품종·체중·알러지가 모든 추천의 출발점",
        "근거 있는 추천 — \"닭고기 알러지 소형견 보호자 12명의 후기를 찾아보니\"처럼 원문 근거를 남김",
        "체크박스로 못 담는 이야기까지 — 정형 데이터와 비정형 자연어를 함께 입력받음",
    ], size=17, color=WHITE, gap=14)
    add_text(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
              "감사합니다 · 오늘뭐먹냥", size=16, color=WHITE, bold=True)

    out = "docs/deck/오늘뭐먹냥_포트폴리오.pptx"
    prs.save(out)
    print(f"saved: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
